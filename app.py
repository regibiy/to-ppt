import os
from flask import Flask, redirect, url_for, render_template, request, send_file
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import ImageFont, Image, ImageOps
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from openpyxl import load_workbook
from openpyxl import get_column_letter
from copy import copy


app = Flask(__name__)

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/ppmreport")
def ppmReport():
    return render_template("to-excel.html")

@app.route("/generate", methods=["POST"]) #pptwarranty
def generate():
    fields = ["caseNumber", "unitModel", "serialNumber", "claimNumber"]
    data = {key: request.form[key].upper() for key in fields}

    prs = Presentation("template.pptx")
    slide = prs.slides[0]

    table = slide.shapes[2].table
    for i, key in enumerate(fields):
        cell = table.cell(i, 1)
        cell.text = data[key]
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = cell.text
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True

    def getPos(shape):
        return {
            "left": shape.left + Inches(0.03),
            "top": shape.top + Inches(0.03),
            "width": shape.width - Inches(0.05),
            "height": shape.height - Inches(0.05)
        }
    
    formatMap = {
        ".jpg": "JPEG",
        ".jpeg": "JPEG",
        ".png": "PNG",
        ".webp": "WEBP"
        }
    
    def processImage(file):
        ext = os.path.splitext(file.filename)[1].lower()
        imgFormat = formatMap.get(ext, "JPEG")

        stream = BytesIO(file.read())
        stream.seek(0)

        img = Image.open(stream)
        img = ImageOps.exif_transpose(img)

        out = BytesIO()
        img.save(out, format=imgFormat)
        out.seek(0)
        return out

    woPos = getPos(slide.shapes[0])
    woStream = processImage(request.files["WO"])

    slide.shapes.add_picture(woStream, woPos["left"], woPos["top"], woPos["width"], woPos["height"])

    # begin get all unit photo 
    photoUnitList = request.files.getlist("addPhoto[]")
    infoUnitList = request.form.getlist("addInfo[]")

    if photoUnitList:
        firstPhoto = photoUnitList[0]
        firstInfo = infoUnitList[0] if infoUnitList else ""

        unitPos = getPos(slide.shapes[1])
        firstStream = processImage(firstPhoto)

        slide.shapes.add_picture(firstStream, unitPos["left"], unitPos["top"], unitPos["width"], unitPos["height"])

        oldInfoPos = getPos(slide.shapes[3])
        slide.shapes._spTree.remove(slide.shapes[3]._element)

        if firstInfo.strip():
            fontPath = os.path.join(os.path.dirname(__file__), "fonts", "calibri-regular.ttf")
            font = ImageFont.truetype(fontPath, 18)
            bbox = font.getbbox(firstInfo)
            textWidthIn = (bbox[2] - bbox[0]) / 96
            newLeft = oldInfoPos["left"] + (oldInfoPos["width"] - Inches(textWidthIn)) / 2

            infoShape = slide.shapes.add_shape(slide.shapes[0].auto_shape_type, newLeft, oldInfoPos["top"], Inches(textWidthIn) + Inches(0.2), oldInfoPos["height"] + Inches(0.2))

            infoShape.text = firstInfo
            infoShape.fill.solid()
            infoShape.fill.fore_color.rgb = RGBColor(0, 176, 80)
            infoShape.line.color.rgb = RGBColor(255, 255, 255)
            infoShape.line.width = Pt(1)

            run = infoShape.text_frame.paragraphs[0].runs[0]
            run.font.color.rgb = RGBColor(255, 255, 255)
    
    slideWidth = prs.slide_width
    margin = Inches(0.3)
    gapY = Inches(1.25)
    infoHeight = Inches(0.5)
    shapeWidth = (slideWidth - (3 * margin)) / 2
    slideLayout = prs.slide_layouts[6]
    slide2 = None

    for i in range(1, len(photoUnitList)):
        photo = photoUnitList[i]
        info = infoUnitList[i] if i < len(infoUnitList) else ""

        if (i - 1) % 4 == 0:
            slide2 = prs.slides.add_slide(slideLayout)

        # index 0–3 dalam slide
        idx = (i - 1) % 4
        row = idx // 2      # 0 atau 1
        col = idx % 2       # 0 atau 1

        x = margin + col * (shapeWidth + margin)
        y = Inches(3.8) if row == 1 else Inches(0.2) + row * (infoHeight + gapY)
        infoBox = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, shapeWidth, infoHeight)

        infoBox.fill.solid()
        infoBox.fill.fore_color.rgb = RGBColor(0, 176, 80)
        infoBox.line.color.rgb = RGBColor(255, 255, 255)
        infoBox.line.width = Pt(1)


        para = infoBox.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.color.rgb = RGBColor(255, 255, 255)
        infoBox.text = info
        
        stream = processImage(photo)
        slide2.shapes.add_picture(stream, x, y + infoHeight, shapeWidth, Inches(3))
          
    pptStream = BytesIO()
    prs.save(pptStream)
    pptStream.seek(0)
    fileName = request.form["caseNumber"] if not request.form["fileName"].strip() else request.form["fileName"].strip()
    return send_file(
        pptStream,
        as_attachment=True,
        download_name = f"{fileName}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.route("/generate2", methods=["POST"]) #excelppm
def generate2():
    def copyRow(ws, srcRow, tgtRow):
        for col in range(1, ws.max_column + 1):
            colLetter = get_column_letter(col)


if __name__ == "__main__":
    app.run(debug=True) 